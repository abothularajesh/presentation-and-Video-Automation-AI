from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import random
import re


def create_ppt(slide_data, filename):

    themes = [
        "templates/Berlin.pptx",
        "templates/Facet.pptx",
        "templates/Celestial.pptx",
        "templates/Depth.pptx",
        "templates/Mesh.pptx",
        "templates/Parallex.pptx",
        "templates/gallery.pptx",
        "templates/ion.pptx",
        "templates/BoardRoom.pptx",
        "templates/Circuit.pptx",
        "templates/Damask.pptx",
        "templates/Droplet.pptx",
        "templates/View.pptx",
        "templates/Slate.pptx",
        "templates/Wisp.pptx"
    ]

    template = random.choice(themes)

    prs = Presentation(template)

    # remove template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


    for i, slide in enumerate(slide_data):

        # remove numbering like "1.", "2)", "3 -"
        title = slide["title"].replace("*", "").strip()
        title = re.sub(r'^\d+[\.\)\-\s]*:', '', title)
        points = slide["points"][:4]
        image = slide.get("image")


        # -------- TITLE SLIDE --------
        if i == 0:

            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)

            clean_title = slide["title"]

            # remove "Here are X slides..." type phrases
            clean_title = re.sub(r"(?i)here are.*?:", "", clean_title)

            # remove numbering if present
            #clean_title = re.sub(r'^\d+[\.\)\-\s]*', '', clean_title)

            # remove words like "explain", "explaining", "about"
            clean_title = re.sub(r'(?i)\b(explain|explaining|about)\b', '', clean_title)

            clean_title = clean_title.strip().title()

            s.shapes.title.text = clean_title
            s.placeholders[1].text = "LearnLift - AI"


       # -------- PROCESS DIAGRAM (SMARTART STYLE) --------
        elif len(points) >= 3 and i % 4 == 0:

            layout = prs.slide_layouts[5]
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = title

            start_left = Inches(0.5)
            box_width = Inches(3.2)      # slightly wider
            box_height = Inches(2)       # slightly taller
            gap = Inches(0.7)
            top = Inches(3)

            theme_colors = [
                MSO_THEME_COLOR.ACCENT_1,
                MSO_THEME_COLOR.ACCENT_2,
                MSO_THEME_COLOR.ACCENT_3,
                MSO_THEME_COLOR.ACCENT_4,
                MSO_THEME_COLOR.ACCENT_5
            ]

            boxes = []

            for k, text in enumerate(points[:3]):

                left = start_left + k * (box_width + gap)

                shape = s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left,
                    top,
                    box_width,
                    box_height
                )

                shape.fill.solid()
                shape.fill.fore_color.theme_color = theme_colors[k % len(theme_colors)]
                shape.line.color.theme_color = MSO_THEME_COLOR.TEXT_1

                tf = shape.text_frame
                tf.clear()

                # enable wrapping
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = text
                p.alignment = PP_ALIGN.CENTER

                # dynamic font sizing
                if len(text) > 140:
                    p.font.size = Pt(14)
                elif len(text) > 90:
                    p.font.size = Pt(16)
                else:
                    p.font.size = Pt(18)

                boxes.append(shape)

        # -------- IMAGE + BULLET CONTENT SLIDE --------
        elif image and i % 3 == 0:

            layout = prs.slide_layouts[5]   # title only
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = title

            # -------- LEFT TEXT BOX --------
            textbox = s.shapes.add_textbox(
                Inches(0.7),     # left margin
                Inches(2),       # top
                Inches(5),       # width
                Inches(3.8)      # height
            )

            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            for j, p in enumerate(points[:4]):

                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = p

                para.level = 0           # this enables bullet automatically
                para.font.size = Pt(24)
                para.space_after = Pt(10)
                para.font.bold = False


            # -------- RIGHT IMAGE --------
            s.shapes.add_picture(
                image,
                Inches(6.3),     # right side
                Inches(2),       # aligned with text
                height=Inches(3.8)
            )


        # -------- TWO CONTENT SLIDE --------
        elif i % 2 == 0:

            layout = prs.slide_layouts[3]
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = title

            left_box = s.shapes.placeholders[1].text_frame
            right_box = s.shapes.placeholders[2].text_frame

            for j, p in enumerate(points):

                if j < 2:
                    para = left_box.paragraphs[0] if j == 0 else left_box.add_paragraph()
                else:
                    para = right_box.paragraphs[0] if j == 2 else right_box.add_paragraph()

                para.text = p
                para.font.size = Pt(22)


        # -------- NORMAL CONTENT --------
        else:

            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = title
            tf = s.shapes.placeholders[1].text_frame

            for j, p in enumerate(points):

                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = p
                para.font.size = Pt(24)


        # title font
        s.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)


    output_path = f"outputs/slides/{filename}.pptx"
    prs.save(output_path)

    return output_path
